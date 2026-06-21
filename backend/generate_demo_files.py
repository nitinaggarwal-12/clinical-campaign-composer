import os
import sys

# Ensure we can import reportlab and pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("python-pptx not found. Installing...")
    os.system("pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
except ImportError:
    print("reportlab not found. Installing...")
    os.system("pip install reportlab")
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

# Save directly inside the frontend folder so the node server serves them!
DEMO_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "frontend", "demo_files")
if not os.path.exists(DEMO_DIR):
    os.makedirs(DEMO_DIR)

def create_pdf():
    pdf_path = os.path.join(DEMO_DIR, "clinical_brief_keynote189.pdf")
    doc = SimpleDocTemplate(pdf_path, pagesize=letter,
                            rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72)
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'TitleStyle',
        parent=styles['Heading1'],
        fontSize=18,
        leading=22,
        textColor='#0D9488',
        spaceAfter=15
    )
    h2_style = ParagraphStyle(
        'H2Style',
        parent=styles['Heading2'],
        fontSize=14,
        leading=18,
        textColor='#1E293B',
        spaceBefore=12,
        spaceAfter=6
    )
    body_style = ParagraphStyle(
        'BodyStyle',
        parent=styles['Normal'],
        fontSize=10,
        leading=14,
        textColor='#334155',
        spaceAfter=10
    )
    
    story = []
    story.append(Paragraph("CLINICAL SUMMARY BRIEF: KEYNOTE-189 TRIAL", title_style))
    story.append(Paragraph("<b>Brand Target:</b> KEYTRUDA (Pembrolizumab)", body_style))
    story.append(Paragraph("<b>Medication:</b> KEYTRUDA", body_style)) # 100% PURE REAL DRUG NAME!
    story.append(Paragraph("<b>Trial Registration:</b> KEYNOTE-189 Phase III Trial (NCT02578680)", body_style))
    story.append(Spacer(1, 10))
    
    story.append(Paragraph("EFFICACY DATA ENDPOINTS", h2_style))
    story.append(Paragraph("In the landmark KEYNOTE-189 Phase III clinical trial, KEYTRUDA combination therapy demonstrated an extraordinary survival benefit. The trial met its primary endpoint of Overall Survival (OS).", body_style))
    story.append(Paragraph("• <b>Key Efficacy Claim:</b> Demonstrated a clinically proven <b>Overall Response Rate (ORR) of 56% at Week 24</b> in the combination group, compared to only 18.9% in patients receiving chemotherapy alone.", body_style))
    story.append(Paragraph("• <b>Hazard Ratio (HR):</b> OS HR was 0.49, representing a highly significant 51% reduction in the risk of death.", body_style))
    story.append(Spacer(1, 10))
    
    story.append(Paragraph("SAFETY & TOLERABILITY PROFILE", h2_style))
    story.append(Paragraph("The safety profile of KEYTRUDA combination therapy was consistent with previously established clinical registries.", body_style))
    story.append(Paragraph("• <b>Key Safety Parameter:</b> Grade 3/4 Immune-Mediated Adverse Reactions occurred in <b>10% of patients</b> in the active cohort.", body_style))
    story.append(Paragraph("• <b>Common Events:</b> Fatigue, nausea, and pneumonitis were monitored according to standard safety guidelines.", body_style))
    story.append(Spacer(1, 10))
    
    story.append(Paragraph("REGULATORY & COMPLIANCE SIGNATURES", h2_style))
    story.append(Paragraph("This brief is approved for clinical marketing campaign generation.", body_style))
    story.append(Paragraph("<b>Regulatory Compliance Vault Ref:</b> Ref #V-2026-KT189-NSCLC", body_style))
    story.append(Paragraph("<b>Trial Registry:</b> KEYNOTE-189 Safety Registry", body_style))
    
    doc.build(story)
    print(f"Created PDF: {pdf_path}")

def create_pptx():
    pptx_path = os.path.join(DEMO_DIR, "clinical_brief_lenvima.pptx")
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "CLINICAL SUMMARY BRIEF: CLEAR TRIAL"
    subtitle.text = "Brand Target: LENVIMA (Lenvatinib)\nMedication: LENVIMA\nTrial Registration: CLEAR Phase III (NCT02811861)" # 100% PURE REAL DRUG NAME!
    
    # Slide 2: Efficacy Endpoints
    slide_layout = prs.slide_layouts[1] # Bullet points
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "EFFICACY DATA ENDPOINTS (CLEAR Trial)"
    tf = body_shape.text_frame
    tf.text = "In the pivotal CLEAR Phase III trial, LENVIMA combination demonstrated a superior efficacy standard:"
    
    p = tf.add_paragraph()
    p.text = "• Key Efficacy Claim: Demonstrated an Objective Response Rate (ORR) of 71% in combination group, compared to 36% with sunitinib alone."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Progression-Free Survival (PFS): Median PFS was 23.9 months versus 9.2 months for sunitinib (highly significant)."
    p.level = 1
    
    # Slide 3: Safety Profile
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    shapes.title.text = "SAFETY & TOLERABILITY PROFILE"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Safety parameters remain consistent with established clinical guidelines:"
    
    p = tf.add_paragraph()
    p.text = "• Key Safety Parameter: Grade 3/4 Adverse Events occurred in 30% of patients in the active cohort."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Adverse Events: Primary presentations included manageable hypertension and fatigue."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Regulatory Vault Ref: Ref #V-2026-LV581-RCC"
    p.level = 1
    
    prs.save(pptx_path)
    print(f"Created PPTX: {pptx_path}")

if __name__ == "__main__":
    create_pdf()
    create_pptx()
    print("All pure clinical demo files successfully generated in frontend/demo_files!")
