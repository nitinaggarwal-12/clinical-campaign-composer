import sys
import os

# Insert backend directory to allow imports
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt

def generate_pptx():
    prs = Presentation()
    
    # Use standard 16:9 widescreen layout
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Drug Indication Profile
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Product-A (compound_alpha) Clinical Indication Profile"
    
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Medication: Product-A (compound_alpha)"
    p.font.size = Pt(22)
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = "Indication: First-line metastatic non-squamous non-small cell lung cancer (NSCLC) with no EGFR or ALK genomic tumor aberrations."
    p2.font.size = Pt(18)
    
    p3 = tf.add_paragraph()
    p3.text = "Target Population: Adults in the first-line setting."
    p3.font.size = Pt(18)
    
    p4 = tf.add_paragraph()
    p4.text = "Core Marketing Hook: Redefining overall survival boundaries with Product-A in first-line NSCLC."
    p4.font.size = Pt(18)
    p4.font.italic = True
    
    # Slide 2: Efficacy Parameters
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "KEYNOTE-189 Trial Efficacy Outcomes"
    
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Clinical Trial: KEYNOTE-189 Phase III Study (NCT02578680)"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = "Key Efficacy Claim: Product-A (compound_alpha) Efficacy: 56% Overall Response Rate (ORR) at Week 24 (KEYNOTE-189 study)."
    p2.font.size = Pt(18)
    
    p3 = tf.add_paragraph()
    p3.text = "Regulatory Compliance Vault Reference: ComplianceVault Ref #V-2026-KT089"
    p3.font.size = Pt(16)
    
    # Slide 3: Safety Parameters
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Immune-Mediated Adverse Reactions & Safety Profile"
    
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Safety Parameter: Grade 3/4 Immune-Mediated Adverse Reactions"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = "Key Safety Claim: Product-A Safety Profile: 10% Grade 3/4 Immune-Mediated Adverse Reactions."
    p2.font.size = Pt(18)
    
    p3 = tf.add_paragraph()
    p3.text = "Regulatory Compliance Vault Reference: ComplianceVault Ref #V-2026-KTS99"
    p3.font.size = Pt(16)
    
    p4 = tf.add_paragraph()
    p4.text = "Regulatory Warning: Severe and fatal immune-mediated adverse reactions can occur in any organ system or tissue."
    p4.font.size = Pt(14)
    p4.font.bold = True
    
    # Save the real PPTX to the project root
    workspace_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    pptx_path = os.path.join(workspace_root, "KEYNOTE-189_Clinical_Launch_Briefing.pptx")
    prs.save(pptx_path)
    print(f"Success: Real PowerPoint file generated at {pptx_path}")

if __name__ == "__main__":
    generate_pptx()
